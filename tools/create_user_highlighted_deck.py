from __future__ import annotations

from pathlib import Path
from textwrap import wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

import create_functionality_testing_deck as base


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "VibeSec_Functionality_Testing_and_UML_Diagrams.pptx"
ASSETS = ROOT / "docs" / "highlighted-function-screens"

W, H = 1366, 768
BG = (7, 14, 10)
PANEL = (18, 27, 20)
PANEL2 = (24, 35, 26)
LINE = (45, 61, 48)
TEXT = (245, 248, 240)
MUTED = (169, 184, 169)
FAINT = (116, 133, 116)
GREEN = (91, 255, 0)
ACCENT = (147, 204, 88)
RED = (248, 88, 88)
YELLOW = (246, 198, 68)
BLUE = (141, 181, 255)


def font(size: int, bold: bool = False, mono: bool = False):
    names = []
    if mono:
        names = [r"C:\Windows\Fonts\consola.ttf"]
    elif bold:
        names = [r"C:\Windows\Fonts\segoeuib.ttf", r"C:\Windows\Fonts\arialbd.ttf"]
    else:
        names = [r"C:\Windows\Fonts\segoeui.ttf", r"C:\Windows\Fonts\arial.ttf"]
    for name in names:
        if Path(name).exists():
            return ImageFont.truetype(name, size)
    return ImageFont.load_default()


F_TITLE = font(24, True)
F_H = font(19, True)
F = font(15)
F_SMALL = font(13)
F_MONO = font(14, mono=True)
F_MONO_SMALL = font(12, mono=True)


def txt(d: ImageDraw.ImageDraw, xy, text: str, fill=TEXT, fnt=F, max_width: int | None = None, spacing: int = 4):
    x, y = xy
    if not max_width:
        d.text((x, y), text, fill=fill, font=fnt)
        return
    chars = max(12, int(max_width / max(7, fnt.size * 0.55)))
    for line in wrap(text, width=chars):
        d.text((x, y), line, fill=fill, font=fnt)
        y += fnt.size + spacing


def rect(d, xy, fill=PANEL, outline=LINE, width=1, radius=8):
    d.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def hi(d, xy, width=6):
    d.rectangle(xy, outline=GREEN, width=width)


def button(d, x, y, w, h, label, fill=(10, 16, 12), outline=LINE, text=TEXT):
    rect(d, (x, y, x + w, y + h), fill=fill, outline=outline, radius=8)
    tw = d.textlength(label, font=F_SMALL)
    d.text((x + (w - tw) / 2, y + 9), label, fill=text, font=F_SMALL)


def side_nav(d, active="Dashboard", control_center=True):
    d.rectangle((0, 0, 255, H), fill=(5, 12, 8))
    rect(d, (14, 24, 50, 60), fill=ACCENT, outline=ACCENT, radius=8)
    txt(d, (58, 24), "VibeSec", fnt=F_H)
    txt(d, (58, 49), "CONTROL CENTER" if control_center else "ANALYSIS", fill=FAINT, fnt=F_MONO_SMALL)
    d.line((0, 86, 245, 86), fill=LINE)
    txt(d, (14, 124), "WORKSPACE", fill=FAINT, fnt=F_MONO_SMALL)
    items = ["Dashboard", "Rules", "Logs", "Settings"] if control_center else ["Analysis"]
    y = 158
    for item in items:
        if item == active:
            rect(d, (0 if control_center else 8, y - 12, 250 if control_center else 245, y + 31), fill=PANEL2, outline=LINE, radius=6)
            fill = TEXT
        else:
            fill = MUTED
        txt(d, (38, y), item, fill=fill, fnt=F_H if item == active else F)
        y += 46
    d.line((0, H - 76, 245, H - 76), fill=LINE)
    txt(d, (12, H - 48), "idle", fill=ACCENT, fnt=F_SMALL)
    txt(d, (18, H - 25), "v1.0.0", fill=FAINT, fnt=F_MONO_SMALL)


def top(d, title, crumb):
    d.rectangle((255, 0, W, 52), fill=(10, 18, 13))
    txt(d, (282, 16), title, fnt=F_H)
    txt(d, (360, 18), f"vibesec > {crumb}", fill=FAINT, fnt=F_MONO_SMALL)
    button(d, W - 155, 8, 135, 36, "Scan project", fill=(13, 21, 15))


def save(path: Path, img: Image.Image):
    ASSETS.mkdir(parents=True, exist_ok=True)
    img.save(path)
    return path


def settings_api_save():
    img = Image.new("RGB", (W, H), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Settings")
    top(d, "Settings", "settings")
    y = 84
    for provider in ["Anthropic", "OpenAI", "Gemini", "Groq", "Custom / Other"]:
        rect(d, (285, y, 1058, y + 104), fill=PANEL, outline=LINE)
        txt(d, (305, y + 20), f"{provider} API key", fnt=F_H)
        if provider == "Groq":
            txt(d, (438, y + 23), "active provider", fill=FAINT, fnt=F_MONO_SMALL)
        txt(d, (305, y + 47), f"{provider} API key. Paste the key, press Save, then use Test to verify it.", fill=MUTED, fnt=F_SMALL, max_width=360)
        txt(d, (305, y + 82), "value: hidden for security", fill=FAINT, fnt=F_MONO_SMALL)
        rect(d, (735, y + 24, 1035, y + 58), fill=(4, 9, 6), outline=LINE, radius=6)
        txt(d, (748, y + 33), f"Paste {provider} key", fill=FAINT, fnt=F_MONO)
        button(d, 830, y + 66, 90, 36, "Save & use", fill=ACCENT, outline=ACCENT, text=(0, 0, 0))
        button(d, 928, y + 66, 54, 36, "Test")
        txt(d, (998, y + 76), "Clear", fill=MUTED, fnt=F_SMALL)
        y += 112
    rect(d, (680, 708, 1080, 756), fill=PANEL2, outline=LINE)
    txt(d, (715, 725), "Groq API key saved and selected as the active provider.", fnt=F_SMALL)
    hi(d, (817, 491, 920, 604))
    hi(d, (675, 706, 1083, 759))
    return save(ASSETS / "01_api_key_save.png", img)


def settings_api_test():
    img = Image.new("RGB", (920, 560), BG)
    d = ImageDraw.Draw(img)
    y = 28
    for provider in ["Gemini", "Groq", "Custom / Other"]:
        rect(d, (70, y, 830, y + 118), fill=PANEL, outline=LINE)
        txt(d, (90, y + 18), f"{provider} API key" + ("  active provider" if provider == "Groq" else ""), fnt=F_H)
        txt(d, (90, y + 48), f"{provider} API key. Paste the key, press Save, then use Test to verify it.", fill=MUTED, fnt=F_SMALL, max_width=380)
        txt(d, (90, y + 92), "value: hidden for security", fill=FAINT, fnt=F_MONO_SMALL)
        rect(d, (510, y + 22, 810, y + 56), fill=(4, 9, 6), outline=LINE, radius=6)
        txt(d, (525, y + 31), f"Paste {provider} key", fill=FAINT, fnt=F_MONO)
        button(d, 604, y + 66, 88, 36, "Save & use", fill=ACCENT, outline=ACCENT, text=(0, 0, 0))
        button(d, 700, y + 66, 50, 36, "Test")
        txt(d, (765, y + 77), "Clear", fill=MUTED, fnt=F_SMALL)
        if provider == "Groq":
            hi(d, (692, y + 62, 759, y + 105))
        y += 136
    rect(d, (684, 486, 842, 534), fill=PANEL2, outline=LINE)
    txt(d, (704, 504), "Groq API key works.", fnt=F_SMALL)
    hi(d, (682, 484, 846, 537))
    return save(ASSETS / "02_api_key_test.png", img)


def logs_page():
    img = Image.new("RGB", (1120, 850), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Logs")
    top(d, "Logs", "logs")
    cards = [("EVENTS", "493", "in buffer"), ("ERRORS", "60", "needs attention"), ("WARNINGS", "59", "non-blocking"), ("INFO", "374", "normal activity")]
    x = 292
    for label, num, sub in cards:
        rect(d, (x, 52, x + 180, 154), fill=PANEL2, outline=LINE)
        txt(d, (x + 18, 66), label, fill=FAINT, fnt=F_MONO_SMALL)
        txt(d, (x + 18, 96), num, fill=RED if label == "ERRORS" else YELLOW if label == "WARNINGS" else TEXT, fnt=font(28, True))
        txt(d, (x + 75, 108), sub, fill=FAINT, fnt=F_MONO_SMALL)
        x += 193
    rect(d, (292, 172, 1052, 850), fill=PANEL, outline=LINE)
    txt(d, (305, 188), "TIME", fill=FAINT, fnt=F_MONO_SMALL)
    txt(d, (458, 188), "TYPE", fill=FAINT, fnt=F_MONO_SMALL)
    txt(d, (586, 188), "LEVEL", fill=FAINT, fnt=F_MONO_SMALL)
    txt(d, (724, 188), "MESSAGE", fill=FAINT, fnt=F_MONO_SMALL)
    rows = [
        ("21:08:52", "API", "INFO", "groq request succeeded (336ms)", "model=llama-3.1-8b-instant\nresponse-length=5"),
        ("21:08:52", "API", "INFO", "groq request started", ""),
        ("21:06:45", "API", "ERROR", "groq request failed (HTTP 413) after 419ms", ""),
        ("21:06:44", "PROMPT", "INFO", "Building per-file prompt - 14 findings", ""),
        ("20:33:52", "SCAN", "INFO", "Scan completed - 14 findings across 1 file", ""),
        ("20:33:52", "SEMGREP", "WARN", "Semgrep wrote to stderr while scanning", ""),
    ]
    y = 222
    for t, typ, level, msg, detail in rows:
        txt(d, (305, y), t, fill=FAINT, fnt=F_MONO_SMALL)
        txt(d, (458, y), typ, fill=MUTED, fnt=F_MONO_SMALL)
        txt(d, (586, y), level, fill=RED if level == "ERROR" else YELLOW if level == "WARN" else MUTED, fnt=F_MONO_SMALL)
        txt(d, (724, y), msg, fnt=F)
        if detail:
            rect(d, (306, y + 34, 1020, y + 100), fill=(5, 10, 7), outline=LINE)
            txt(d, (321, y + 50), detail, fill=MUTED, fnt=F_MONO)
            y += 130
        else:
            y += 42
    hi(d, (288, 48, 1058, 160))
    hi(d, (720, 217, 1030, 340))
    return save(ASSETS / "03_logs_api_success.png", img)


def rules_policy_toggles():
    img = Image.new("RGB", (880, 768), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Rules")
    top(d, "Rules", "rules")
    stats = [("TOTAL RULES", "154", "154 enabled"), ("ACTIVE FILES", "2", "policies"), ("CUSTOM", "1", "policy files")]
    x = 68
    for label, num, sub in stats:
        rect(d, (x, 68, x + 250, 158), fill=PANEL2, outline=ACCENT if label == "TOTAL RULES" else LINE)
        txt(d, (x + 18, 88), label, fill=FAINT, fnt=F_MONO_SMALL)
        txt(d, (x + 18, 116), num, fill=ACCENT if label != "CUSTOM" else (255, 198, 120), fnt=font(27, True))
        txt(d, (x + 68, 127), sub, fill=MUTED, fnt=F_MONO_SMALL)
        x += 264
    txt(d, (67, 205), "Turn ON any number of policy files. Normal and taint policies can run together, or all can be OFF.", fill=MUTED, fnt=F, max_width=560)
    policies = [
        ("rules/default.yaml", "OWASP Top 10 baseline - injection, crypto, secrets, XSS, auth, integrity.", "141", "rules", True),
        ("rules/taint.yaml", "Taint analysis - tracks user input from source to dangerous sink within a file.", "13", "rules", True),
        ("rules/policies/normal-sara.yaml", "Tool policy file - stored inside VibeSec's rules/policies folder.", "0", "rules", False),
    ]
    y = 240
    for name, desc, count, unit, on in policies:
        rect(d, (67, y, 847, y + 125), fill=PANEL, outline=LINE)
        txt(d, (142, y + 22), name, fnt=F_H)
        txt(d, (142, y + 52), desc, fill=MUTED, fnt=F_SMALL, max_width=480)
        rect(d, (672, y + 37, 737, y + 89), fill=(6, 11, 7), outline=LINE)
        txt(d, (690, y + 50), count, fnt=F_H)
        txt(d, (690, y + 74), unit, fill=FAINT, fnt=F_MONO_SMALL)
        d.rounded_rectangle((753, y + 50, 795, y + 75), 13, fill=ACCENT if on else (40, 51, 43))
        d.ellipse((775 if on else 758, y + 53, 792 if on else 775, y + 70), fill=BG)
        y += 132
    hi(d, (704, 243, 831, 506))
    return save(ASSETS / "04_rules_policy_toggles.png", img)


def rules_individual_toggle():
    img = Image.new("RGB", (1318, 912), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Rules")
    top(d, "Rules", "rules")
    names = ["Ssrf axios from user url", "Ssrf fetch from user url", "Ssrf http request from...", "Ssrf image fetch thum...", "Cors wildcard origin e..."]
    y = 90
    for i, name in enumerate(names):
        rect(d, (280, y, 1040, y + 115), fill=PANEL, outline=LINE, radius=0)
        sev = "WARNING" if i == 3 or i == 4 else "ERROR"
        rect(d, (296, y + 24, 383, y + 50), fill=(45, 20, 20) if sev == "ERROR" else (45, 37, 10), outline=RED if sev == "ERROR" else YELLOW)
        txt(d, (312, y + 29), sev, fill=RED if sev == "ERROR" else YELLOW, fnt=F_MONO_SMALL)
        txt(d, (398, y + 24), name, fnt=F_H)
        txt(d, (398, y + 52), "vibesec.ssrf-rule-id", fill=FAINT, fnt=F_MONO_SMALL)
        txt(d, (584, y + 42), "security", fnt=F_SMALL)
        txt(d, (724, y + 42), "CWE-918", fill=MUTED, fnt=F_SMALL)
        d.rounded_rectangle((980, y + 38, 1020, y + 62), 12, fill=ACCENT if i < 4 else (40, 51, 43))
        d.ellipse((1000 if i < 4 else 985, y + 41, 1017 if i < 4 else 1002, y + 58), fill=BG)
        y += 112
    rect(d, (890, 842, 1310, 898), fill=PANEL2, outline=LINE)
    txt(d, (929, 861), "VibeSec: Policy reloaded successfully.", fnt=F)
    hi(d, (912, 605, 1051, 688))
    hi(d, (931, 747, 1068, 834))
    hi(d, (888, 840, 1312, 900))
    return save(ASSETS / "05_individual_rule_toggle.png", img)


def analysis_generate():
    img = Image.new("RGB", (1912, 872), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Analysis", control_center=False)
    txt(d, (70, 36), "Analysis", fnt=F_TITLE)
    txt(d, (70, 70), "Scan source for vulnerabilities and generate ready-to-paste fix prompts.", fill=MUTED, fnt=F)
    rect(d, (70, 164, 396, 340), fill=PANEL, outline=LINE)
    txt(d, (69, 144), "FILES 1", fill=FAINT, fnt=F_MONO_SMALL)
    txt(d, (120, 180), "1 of 4 files selected", fill=MUTED, fnt=F_SMALL)
    txt(d, (98, 254), "insecure.py", fnt=F_H)
    button(d, 70, 353, 326, 41, "Analyze", fill=ACCENT, outline=ACCENT, text=(0, 0, 0))
    txt(d, (70, 423), "OUTPUT", fill=FAINT, fnt=F_MONO_SMALL)
    button(d, 31, 440, 166, 36, "Results  9")
    button(d, 194, 440, 156, 36, "Full Fix  1")
    rect(d, (62, 492, 402, 573), fill=PANEL, outline=LINE)
    txt(d, (92, 508), "9 fix\nprompts\nacross 1 file", fnt=F_H)
    button(d, 185, 518, 100, 31, "Generate", fill=(45, 72, 32), outline=ACCENT, text=ACCENT)
    button(d, 298, 518, 86, 31, "Copy all")
    rect(d, (1310, 687, 1910, 797), fill=PANEL2, outline=LINE)
    txt(d, (1350, 720), "VibeSec: Generating prompts (perFile, Groq)... File 1/1 - insecure.py", fnt=F)
    txt(d, (1350, 762), "Source: VibeSec", fill=MUTED, fnt=F_SMALL)
    hi(d, (172, 507, 303, 564))
    hi(d, (1310, 686, 1911, 798))
    return save(ASSETS / "06_generate_prompt_progress.png", img)


def results_findings():
    img = Image.new("RGB", (435, 910), BG)
    d = ImageDraw.Draw(img)
    txt(d, (33, 6), "Analysis", fnt=F_H)
    button(d, 38, 136, 132, 40, "Results  9")
    txt(d, (48, 187), "All 9", fill=ACCENT, fnt=F)
    rect(d, (32, 226, 370, 510), fill=PANEL, outline=(92, 41, 41))
    txt(d, (49, 247), "ERROR", fill=RED, fnt=F_MONO_SMALL)
    txt(d, (49, 272), "vibesec.command-injection-os-system", fill=FAINT, fnt=F_MONO_SMALL)
    txt(d, (49, 314), "Command injection os system", fnt=F_H)
    txt(d, (49, 342), "os.system() passes its argument directly to the system shell. Replace with subprocess.run([...], shell=False) so arguments are passed as a list and metacharacters lose their special meaning.", fill=MUTED, fnt=F_SMALL, max_width=285)
    rect(d, (32, 521, 371, 814), fill=PANEL, outline=(92, 41, 41))
    txt(d, (49, 544), "ERROR", fill=RED, fnt=F_MONO_SMALL)
    txt(d, (49, 570), "vibesec.weak-hash-md5-python", fill=FAINT, fnt=F_MONO_SMALL)
    txt(d, (49, 592), "Weak hash md5 python", fnt=F_H)
    txt(d, (49, 621), "MD5 is broken for any security purpose. For password hashing use bcrypt, scrypt, or argon2; for general-purpose hashing use SHA-256 or SHA-3.", fill=MUTED, fnt=F_SMALL, max_width=285)
    hi(d, (33, 132, 173, 179))
    hi(d, (34, 291, 283, 343))
    hi(d, (30, 521, 243, 613))
    return save(ASSETS / "07_results_findings.png", img)


def dashboard_actions():
    img = Image.new("RGB", (1100, 868), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Dashboard")
    top(d, "Dashboard", "dashboard")
    txt(d, (308, 58), "267", fnt=font(31, True))
    txt(d, (366, 68), "findings", fill=MUTED, fnt=F)
    rect(d, (288, 230, 682, 410), fill=PANEL, outline=LINE)
    for i, (sev, num, color) in enumerate([("ERROR", "234", RED), ("WARNING", "33", YELLOW), ("INFO", "0", BLUE)]):
        x = 305 + i * 126
        rect(d, (x, 280, x + 112, 402), fill=PANEL2, outline=color)
        txt(d, (x + 16, 302), sev, fill=color, fnt=F_MONO_SMALL)
        txt(d, (x + 16, 336), num, fnt=font(30, True))
    rect(d, (678, 32, 1072, 462), fill=PANEL, outline=LINE)
    txt(d, (698, 52), "QUICK ACTIONS", fill=FAINT, fnt=F_MONO_SMALL)
    actions = ["Scan project\nRun full Semgrep sweep", "Open policy file\nChoose any VibeSec policy to open", "Reload policy\nRe-parse rules from disk", "New normal policy\nCreate a named normal scan policy", "New taint policy\nCreate a named taint scan policy"]
    y = 82
    for a in actions:
        rect(d, (698, y, 1055, y + 65), fill=PANEL2, outline=LINE)
        txt(d, (764, y + 17), a, fnt=F_H if "\n" in a else F)
        y += 75
    rect(d, (678, 711, 1073, 867), fill=PANEL, outline=LINE)
    txt(d, (698, 742), "TREND - LAST 7 DAYS", fill=FAINT, fnt=F_MONO_SMALL)
    d.line((706, 776, 746, 824), fill=ACCENT, width=3)
    d.line((746, 824, 982, 825), fill=ACCENT, width=3)
    hi(d, (271, 267, 682, 412))
    hi(d, (708, 84, 929, 153))
    hi(d, (708, 160, 943, 210))
    hi(d, (701, 236, 897, 286))
    hi(d, (713, 312, 1032, 377))
    hi(d, (707, 386, 1025, 444))
    hi(d, (677, 710, 1074, 868))
    return save(ASSETS / "08_dashboard_actions.png", img)


def open_policy():
    img = Image.new("RGB", (1173, 778), BG)
    d = ImageDraw.Draw(img)
    dashboard_actions_bg(d)
    rect(d, (20, 6, 765, 310), fill=(30, 32, 34), outline=LINE)
    txt(d, (291, 18), "VibeSec - Open policy file", fnt=F)
    rect(d, (22, 46, 755, 80), fill=(20, 25, 27), outline=(80, 140, 170))
    options = [
        ("Bundled normal scan policy", "rules/default.yaml", "Default VibeSec rules shipped with the extension"),
        ("Bundled taint policy", "rules/taint.yaml", "Taint source-to-sink rules shipped with the extension"),
        ("normal-sara.yaml", "tool policy folder", ""),
        (".vibesec.yaml", "workspace selector", ""),
    ]
    y = 88
    for i, (a, b, c) in enumerate(options):
        if i == 0:
            d.rectangle((24, y, 756, y + 56), fill=(44, 132, 165))
            fill = TEXT
        else:
            fill = MUTED
        txt(d, (34, y + 10), f"{a}  {b}", fill=fill, fnt=F)
        if c:
            txt(d, (34, y + 36), c, fill=(210, 220, 220), fnt=F_SMALL)
        y += 56
    hi(d, (2, 69, 763, 309))
    hi(d, (795, 243, 1094, 323))
    return save(ASSETS / "09_open_policy_file.png", img)


def dashboard_actions_bg(d):
    side_nav(d, "Dashboard")
    top(d, "Dashboard", "dashboard")
    rect(d, (768, 180, 1150, 563), fill=PANEL, outline=LINE)
    actions = ["Scan project", "Open policy file", "Reload policy", "New normal policy", "New taint policy"]
    y = 180
    for action in actions:
        rect(d, (790, y, 1145, y + 64), fill=PANEL2, outline=LINE)
        txt(d, (855, y + 18), action, fnt=F_H)
        y += 75
    txt(d, (400, 156), "267", fnt=font(31, True))
    txt(d, (442, 168), "findings", fill=MUTED, fnt=F)


def new_policy(kind="normal"):
    img = Image.new("RGB", (1191, 623), BG)
    d = ImageDraw.Draw(img)
    dashboard_actions_bg(d)
    title = "VibeSec - New normal policy" if kind == "normal" else "VibeSec - New taint policy"
    value = "normal-baseline" if kind == "normal" else "taint-api-checks"
    button_text = "New normal policy" if kind == "normal" else "New taint policy"
    rect(d, (8, 6, 785, 141), fill=(29, 31, 32), outline=LINE)
    txt(d, (290, 18), title, fnt=F)
    rect(d, (24, 47, 740, 81), fill=(17, 23, 24), outline=(80, 150, 180))
    txt(d, (30, 58), value, fill=FAINT, fnt=F_MONO)
    txt(d, (29, 92), "Enter a name. VibeSec will create a separate YAML file inside the tool policy folder: rules/policies/.", fnt=F)
    txt(d, (29, 116), "Press 'Enter' to confirm or 'Escape' to cancel", fnt=F)
    y = 405 if kind == "normal" else 480
    hi(d, (8, 6, 785, 141))
    hi(d, (820, y, 1122, y + 69))
    return save(ASSETS / f"10_new_{kind}_policy.png", img)


def copy_all():
    img = Image.new("RGB", (1811, 766), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Analysis", control_center=False)
    txt(d, (28, 34), "FILES 1", fill=FAINT, fnt=F_MONO_SMALL)
    rect(d, (28, 65, 356, 231), fill=PANEL, outline=LINE)
    txt(d, (69, 76), "1 of 4 files selected", fill=MUTED, fnt=F_SMALL)
    txt(d, (86, 146), "insecure.py", fnt=F_H)
    button(d, 28, 245, 328, 42, "Analyze", fill=ACCENT, outline=ACCENT, text=(0, 0, 0))
    txt(d, (28, 318), "OUTPUT", fill=FAINT, fnt=F_MONO_SMALL)
    button(d, 29, 340, 156, 36, "Results  9")
    button(d, 194, 340, 156, 36, "Full Fix  1")
    rect(d, (22, 461, 403, 541), fill=PANEL, outline=LINE)
    txt(d, (52, 488), "9 fix\nprompts\nacross 1 file", fnt=F_H)
    button(d, 144, 411, 100, 31, "Generate", fill=(45, 72, 32), outline=ACCENT, text=ACCENT)
    button(d, 264, 410, 118, 42, "Copy all")
    rect(d, (1282, 697, 1720, 758), fill=PANEL2, outline=LINE)
    txt(d, (1318, 724), "VibeSec: Project prompt copied to clipboard.", fnt=F)
    hi(d, (263, 398, 383, 453))
    hi(d, (1282, 697, 1720, 758))
    return save(ASSETS / "12_copy_all_prompt.png", img)


def scan_four_files():
    img = Image.new("RGB", (1842, 831), BG)
    d = ImageDraw.Draw(img)
    side_nav(d, "Analysis", control_center=False)
    txt(d, (2, 37), "FILES 4", fill=FAINT, fnt=F_MONO_SMALL)
    rect(d, (1, 134, 328, 309), fill=PANEL, outline=LINE)
    txt(d, (53, 147), "4 of 4 files selected", fill=MUTED, fnt=F_SMALL)
    for i, name in enumerate([".vibesec.yaml", "insecure.py", "simple_web_policy_test.js", "vibesec_11_normal_3_taint..."]):
        txt(d, (69, 184 + i * 33), name, fnt=F_H)
    button(d, 1, 322, 328, 42, "Analyze", fill=ACCENT, outline=ACCENT, text=(0, 0, 0))
    rect(d, (660, 230, 1028, 353), fill=PANEL, outline=LINE)
    for i, (sev, num, color) in enumerate([("ERROR", "234", RED), ("WARNING", "33", YELLOW), ("INFO", "0", BLUE)]):
        x = 680 + i * 126
        rect(d, (x, 230, x + 112, 353), fill=PANEL2, outline=color)
        txt(d, (x + 18, 254), sev, fill=color, fnt=F_MONO_SMALL)
        txt(d, (x + 18, 288), num, fnt=font(30, True))
    rect(d, (1240, 704, 1840, 824), fill=PANEL2, outline=LINE)
    txt(d, (1282, 752), "VibeSec: Scanning 4 files... (1/4) .vibesec.yaml", fnt=F)
    button(d, 1755, 785, 64, 32, "Cancel", fill=(67, 145, 190))
    hi(d, (1069, 27, 1316, 103))
    hi(d, (1239, 704, 1841, 831))
    return save(ASSETS / "13_scan_four_files.png", img)


def make_deck():
    ASSETS.mkdir(parents=True, exist_ok=True)
    slides = [
        ("API Key Save & Active Provider", "Settings screen proves Save & use stores Groq key and selects the provider.", settings_api_save()),
        ("API Key Test", "Test button verifies the saved Groq key and shows success toast.", settings_api_test()),
        ("Logs Evidence", "Logs page records API, prompt, scan, semgrep, policy, warning, and error activity.", logs_page()),
        ("Policy File Toggles", "Rules page proves default and taint policies can run together, while custom policies can stay off.", rules_policy_toggles()),
        ("Individual Rule Toggle", "Rules list proves single rules can be enabled/disabled and policy reloads.", rules_individual_toggle()),
        ("Generate Fix Prompts", "Full Fix tab triggers Groq per-file prompt generation with progress notification.", analysis_generate()),
        ("Finding Results", "Results tab shows concrete vulnerabilities and counts from the running scan.", results_findings()),
        ("Dashboard Actions", "Dashboard proves scan project, open policy, reload policy, new policy actions, severity stats, and trend.", dashboard_actions()),
        ("Open Policy File", "Open policy function shows bundled, tool-folder, and workspace policy choices.", open_policy()),
        ("New Normal Policy", "Quick action opens a naming prompt for a new normal policy YAML file.", new_policy("normal")),
        ("New Taint Policy", "Quick action opens a naming prompt for a new taint policy YAML file.", new_policy("taint")),
        ("Copy All Prompts", "Copy all function sends the project prompt to the clipboard.", copy_all()),
        ("Scan Multiple Files", "Analyze/Scan project function runs over selected files and shows progress with cancel.", scan_four_files()),
    ]

    prs = Presentation()
    prs.slide_width = Inches(base.SLIDE_W)
    prs.slide_height = Inches(base.SLIDE_H)
    title = prs.slides.add_slide(prs.slide_layouts[6])
    title.background.fill.solid()
    title.background.fill.fore_color.rgb = base.NAVY
    base.add_text(title, 0.8, 1.0, 11.7, 0.6, "VibeSec Functionality Testing", size=36, color=base.WHITE, bold=True)
    base.add_text(title, 0.82, 1.8, 11.6, 0.8, "PowerPoint rebuilt around the highlighted running-tool screenshots", size=25, color=base.WHITE, bold=True)
    base.add_text(title, 0.85, 3.05, 11.0, 0.7, "Covers API keys, logs, rules, policy creation, scanning, findings, prompt generation, clipboard copy, and architecture diagrams.", size=16, color=base.LINE)

    overview = base.new_slide(prs, "Function Coverage From Screenshots", "Each slide maps one highlighted screenshot to the function being tested.")
    rows = [["Slide", "Functionality", "Evidence shown"]]
    for i, (t, sub, _path) in enumerate(slides, start=1):
        rows.append([str(i), t, sub])
    table = overview.shapes.add_table(len(rows), 3, Inches(0.45), Inches(1.18), Inches(12.45), Inches(5.9)).table
    table.columns[0].width = Inches(0.65)
    table.columns[1].width = Inches(2.7)
    table.columns[2].width = Inches(9.1)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.size = base.Pt(7.7 if r else 9)
            p.font.bold = r == 0
            p.font.color.rgb = base.WHITE if r == 0 else base.INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = base.NAVY if r == 0 else base.WHITE

    for t, sub, img_path in slides:
        slide = base.new_slide(prs, t, sub)
        slide.shapes.add_picture(str(img_path), Inches(0.55), Inches(1.18), width=Inches(12.25))

    base.sequence_diagram_slide(prs)
    base.use_case_diagram_slide(prs)
    base.class_diagram_slide(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    make_deck()
