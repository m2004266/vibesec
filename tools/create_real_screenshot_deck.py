from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

import create_functionality_testing_deck as base


ROOT = Path(__file__).resolve().parents[1]
REAL = ROOT / "docs" / "real-screenshots"
OUT = ROOT / "docs" / "VibeSec_Functionality_Testing_and_UML_Diagrams.pptx"
REPORT = ROOT / "docs" / "real-screenshots-report.json"


def add_real_image_slide(prs: Presentation, title: str, subtitle: str, image_name: str) -> None:
    slide = base.new_slide(prs, title, subtitle)
    image = REAL / image_name
    slide.shapes.add_picture(str(image), Inches(0.48), Inches(1.18), width=Inches(12.38))


def real_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = base.NAVY
    base.add_text(slide, 0.82, 1.0, 11.7, 0.55, "VibeSec", size=38, color=base.WHITE, bold=True)
    base.add_text(
        slide,
        0.85,
        1.75,
        11.7,
        1.0,
        "Real Functionality Testing Screenshots and UML Diagrams",
        size=29,
        color=RGBColor(224, 242, 254),
        bold=True,
    )
    base.add_text(
        slide,
        0.88,
        3.0,
        11.2,
        0.8,
        "This deck uses actual screenshots from the VibeSec VS Code Extension Development Host while the extension is running.",
        size=16,
        color=RGBColor(203, 213, 225),
    )
    base.add_text(slide, 0.9, 6.72, 10.8, 0.28, f"Workspace: {ROOT}", size=8.5, color=RGBColor(148, 163, 184))


def real_summary_slide(prs: Presentation) -> None:
    slide = base.new_slide(prs, "Real Screenshot Evidence", "Captured from the running VS Code Extension Development Host.")
    cards = [
        ("Extension Host", "VibeSec activity bar and Analysis webview loaded", "00_vscode_dev_host.png", base.BLUE),
        ("Command Surface", "Command Palette lists VibeSec commands", "15_command_palette_vibesec.png", base.CYAN),
        ("Scan Results", "14 findings shown from real scan command", "10_scan_current_file_command_results.png", base.GREEN),
        ("Control Center", "Dashboard and settings webview running", "11_control_center_dashboard.png", base.PURPLE),
    ]
    x = 0.72
    for title, desc, image, color in cards:
        base.add_rect(slide, x, 1.35, 2.85, 1.35, f"{title}\n{desc}", fill=base.WHITE, line=color, size=10.5, bold=False)
        base.add_text(slide, x + 0.12, 2.82, 2.6, 0.25, image, size=7.5, color=base.MUTED, align=PP_ALIGN.CENTER)
        x += 3.05
    base.add_text(
        slide,
        0.82,
        3.55,
        11.65,
        1.0,
        "Important: these are not generated mockups. They were captured from the running extension window after launching VS Code with --extensionDevelopmentPath, opening the sample project, and executing VibeSec commands.",
        size=16,
        color=base.INK,
    )
    base.add_text(
        slide,
        0.82,
        4.75,
        11.65,
        0.9,
        "The deck still includes the Sequence Diagram, Use Case Diagram, and Class Diagram in the same PowerPoint.",
        size=15,
        color=base.MUTED,
    )


def function_coverage_slide(prs: Presentation) -> None:
    slide = base.new_slide(prs, "Function Coverage Map", "Each VibeSec function is tied to a real running screenshot.")
    rows = [
        ["Function", "Evidence screenshot", "What it proves"],
        ["scanCurrentFile", "Scan results", "Current sample file scanned and 14 findings displayed."],
        ["scanSelected", "Analysis panel", "Selectable file-tree workflow is visible in the running panel."],
        ["scanWorkspace", "Command Palette / Dashboard", "Whole-project scan command and dashboard scan action are visible."],
        ["goToFinding", "Scan results", "Editor diagnostics and finding cards point into the active source file."],
        ["reloadPolicy", "Command Palette / Dashboard", "Reload Policy command and quick action are visible."],
        ["openPolicyFile", "Command Palette", "Open Policy File command is listed in VS Code."],
        ["copyDescription", "Finding card", "Finding card has copy/control icons in the result UI."],
        ["setApiKey", "Command Palette / Settings", "Set API Key command and provider settings are visible."],
        ["clearApiKey", "Command Palette", "Clear API Key command is visible."],
        ["testApiKey", "Command Palette", "Test API Key command is visible."],
        ["generatePrompts", "Command Palette / Analysis panel", "Generate Prompts command and Full Fix tab are visible."],
        ["copyPromptForVuln", "Finding card", "Per-finding prompt/copy surface is visible."],
        ["copyPromptForFile", "Full Fix tab", "File-level fix-prompt tab is visible."],
        ["copyPromptForAll", "Command Palette", "Copy Prompt for All Findings command is visible."],
        ["openControlCenter", "Control Center", "Control Center opens and renders dashboard/settings."],
    ]
    table = slide.shapes.add_table(len(rows), 3, Inches(0.45), Inches(1.15), Inches(12.45), Inches(5.92)).table
    table.columns[0].width = Inches(2.25)
    table.columns[1].width = Inches(3.05)
    table.columns[2].width = Inches(7.15)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.size = Pt(7.6 if r else 9)
            p.font.bold = r == 0
            p.font.color.rgb = base.WHITE if r == 0 else base.INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = base.NAVY if r == 0 else base.WHITE


def notes_slide(prs: Presentation) -> None:
    slide = base.new_slide(prs, "Testing Notes", "Real UI capture plus repeatable local checks.")
    base.add_rect(slide, 0.75, 1.35, 5.75, 4.9, "", fill=base.WHITE, line=base.LINE)
    base.add_text(slide, 1.02, 1.62, 5.1, 0.3, "Captured live", size=16, color=base.GREEN, bold=True)
    live = [
        "VS Code Extension Development Host running VibeSec.",
        "Analysis webview loaded from the extension.",
        "VibeSec command palette entries visible.",
        "Current-file scan executed against sample JavaScript.",
        "Results panel shows 14 findings and editor diagnostics.",
        "Control Center dashboard/settings webview rendered.",
    ]
    y = 2.05
    for item in live:
        base.add_text(slide, 1.02, y, 5.0, 0.28, "- " + item, size=11, color=base.INK)
        y += 0.43

    base.add_rect(slide, 6.82, 1.35, 5.75, 4.9, "", fill=base.WHITE, line=base.LINE)
    base.add_text(slide, 7.08, 1.62, 5.1, 0.3, "Repeatable checks", size=16, color=base.BLUE, bold=True)
    checks = [
        "npm run compile succeeded before capture.",
        "Semgrep is installed and used by the extension.",
        "Rules index contains bundled normal and taint rules.",
        "Generated screenshots are stored under docs/real-screenshots.",
        "This deck is regenerated from those real images.",
    ]
    y = 2.05
    for item in checks:
        base.add_text(slide, 7.08, y, 5.0, 0.28, "- " + item, size=11, color=base.INK)
        y += 0.43


def main() -> None:
    required = [
        "00_vscode_dev_host.png",
        "08_sample_file_open.png",
        "10_scan_current_file_command_results.png",
        "11_control_center_dashboard.png",
        "14_control_center_settings.png",
        "15_command_palette_vibesec.png",
        "16_analysis_full_fix_tab.png",
    ]
    missing = [name for name in required if not (REAL / name).exists()]
    if missing:
        raise SystemExit("Missing real screenshots: " + ", ".join(missing))

    prs = Presentation()
    prs.slide_width = Inches(base.SLIDE_W)
    prs.slide_height = Inches(base.SLIDE_H)

    real_title_slide(prs)
    real_summary_slide(prs)
    function_coverage_slide(prs)
    add_real_image_slide(prs, "Real Screenshot: Extension Host", "VibeSec loaded in VS Code Extension Development Host.", "00_vscode_dev_host.png")
    add_real_image_slide(prs, "Real Screenshot: Command Palette", "VibeSec command functions visible in VS Code.", "15_command_palette_vibesec.png")
    add_real_image_slide(prs, "Real Screenshot: Sample File Open", "Vulnerable JavaScript sample open in the editor.", "08_sample_file_open.png")
    add_real_image_slide(prs, "Real Screenshot: Scan Results", "VibeSec Scan Current File executed; 14 findings shown.", "10_scan_current_file_command_results.png")
    add_real_image_slide(prs, "Real Screenshot: Finding Actions", "Result card, details area, and fix-prompt surface visible.", "16_analysis_full_fix_tab.png")
    add_real_image_slide(prs, "Real Screenshot: Control Center Dashboard", "Dashboard with scan history, severity breakdown, and quick actions.", "11_control_center_dashboard.png")
    add_real_image_slide(prs, "Real Screenshot: Settings/API Surface", "Engine, file-extension, provider, model, and API-related settings.", "14_control_center_settings.png")
    base.sequence_diagram_slide(prs)
    base.use_case_diagram_slide(prs)
    base.class_diagram_slide(prs)
    notes_slide(prs)

    prs.save(OUT)
    REPORT.write_text(
        json.dumps(
            {
                "deck": str(OUT),
                "screenshots": [str(REAL / name) for name in required],
                "note": "Deck evidence slides use real screenshots from the running VS Code extension.",
            },
            indent=2,
        ),
        encoding="utf-8",
    )
    print(OUT)


if __name__ == "__main__":
    main()
