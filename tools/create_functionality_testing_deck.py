from __future__ import annotations

import json
import os
import re
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path
from textwrap import shorten, wrap

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
ASSETS = DOCS / "ppt-assets"
PPTX_OUT = DOCS / "VibeSec_Functionality_Testing_and_UML_Diagrams.pptx"
REPORT_OUT = ASSETS / "functionality-test-report.json"

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(16, 24, 39)
INK = RGBColor(24, 33, 48)
MUTED = RGBColor(91, 105, 128)
BLUE = RGBColor(37, 99, 235)
CYAN = RGBColor(8, 145, 178)
GREEN = RGBColor(22, 163, 74)
AMBER = RGBColor(217, 119, 6)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(124, 58, 237)
LINE = RGBColor(203, 213, 225)
PAPER = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)


@dataclass
class CheckResult:
    name: str
    command: str
    ok: bool
    summary: list[str]
    raw_tail: str = ""
    image: str | None = None


def run_cmd(label: str, args: list[str], timeout: int = 180) -> subprocess.CompletedProcess[str]:
    try:
        return subprocess.run(
            args,
            cwd=ROOT,
            text=True,
            encoding="utf-8",
            errors="replace",
            capture_output=True,
            timeout=timeout,
        )
    except Exception as exc:
        return subprocess.CompletedProcess(args=args, returncode=99, stdout="", stderr=str(exc))


def font(size: int, mono: bool = False, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    candidates = []
    if mono:
        candidates.extend(
            [
                r"C:\Windows\Fonts\consola.ttf",
                r"C:\Windows\Fonts\consolab.ttf" if bold else r"C:\Windows\Fonts\consola.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf",
            ]
        )
    else:
        candidates.extend(
            [
                r"C:\Windows\Fonts\segoeuib.ttf" if bold else r"C:\Windows\Fonts\segoeui.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
            ]
        )
    for item in candidates:
        if item and Path(item).exists():
            return ImageFont.truetype(item, size)
    return ImageFont.load_default()


def text_lines(text: str, width: int = 120, max_lines: int = 25) -> list[str]:
    lines: list[str] = []
    for raw in text.replace("\r\n", "\n").split("\n"):
        if not raw:
            lines.append("")
            continue
        chunks = wrap(raw, width=width, replace_whitespace=False, drop_whitespace=False)
        lines.extend(chunks or [""])
    if len(lines) > max_lines:
        return lines[: max_lines - 1] + [f"... ({len(lines) - max_lines + 1} more lines)"]
    return lines


def screenshot(name: str, title: str, command: str, status: str, body: list[str], raw_tail: str = "") -> str:
    ASSETS.mkdir(parents=True, exist_ok=True)
    path = ASSETS / f"{name}.png"
    width, height = 1600, 900
    img = Image.new("RGB", (width, height), (15, 23, 42))
    d = ImageDraw.Draw(img)

    d.rounded_rectangle((34, 30, width - 34, height - 30), radius=28, fill=(2, 6, 23), outline=(51, 65, 85), width=2)
    d.rounded_rectangle((34, 30, width - 34, 100), radius=28, fill=(30, 41, 59))
    d.rectangle((34, 72, width - 34, 100), fill=(30, 41, 59))
    for i, color in enumerate([(248, 113, 113), (251, 191, 36), (74, 222, 128)]):
        d.ellipse((64 + i * 32, 56, 82 + i * 32, 74), fill=color)
    d.text((170, 51), title, fill=(226, 232, 240), font=font(28, bold=True))

    y = 128
    d.text((64, y), "$ " + command, fill=(125, 211, 252), font=font(24, mono=True))
    y += 48
    status_color = (74, 222, 128) if status.upper().startswith("PASS") else (251, 191, 36)
    d.rounded_rectangle((64, y, 245, y + 44), radius=18, fill=(20, 83, 45) if status_color[1] > 200 else (120, 53, 15))
    d.text((88, y + 8), status, fill=status_color, font=font(22, bold=True))
    y += 76

    for line in body:
        if y > height - 90:
            break
        d.text((76, y), line, fill=(226, 232, 240), font=font(24, mono=True))
        y += 34

    if raw_tail:
        y += 12
        d.line((64, y, width - 64, y), fill=(51, 65, 85), width=2)
        y += 22
        for line in text_lines(raw_tail, width=125, max_lines=8):
            if y > height - 70:
                break
            d.text((76, y), line, fill=(148, 163, 184), font=font(20, mono=True))
            y += 28

    img.save(path)
    return str(path)


def tail(stdout: str, stderr: str, limit: int = 1600) -> str:
    text = (stdout + "\n" + stderr).strip()
    if len(text) <= limit:
        return text
    return text[-limit:]


def build_compile_check() -> CheckResult:
    npm = shutil.which("npm.cmd") or shutil.which("npm") or "npm"
    result = run_cmd("compile", [npm, "run", "compile"], timeout=240)
    ok = result.returncode == 0
    body = [
        f"exit_code: {result.returncode}",
        "TypeScript compile: " + ("passed" if ok else "failed"),
        "Design bundle: media/design/main.js + controlCenter.js",
        "Output entry: out/extension.js",
    ]
    img = screenshot(
        "01_build_compile",
        "Build and Bundle Verification",
        "npm run compile",
        "PASS" if ok else "WARN",
        body,
        tail(result.stdout, result.stderr),
    )
    return CheckResult("Build and bundle", "npm run compile", ok, body, tail(result.stdout, result.stderr), img)


def build_command_registration_check() -> CheckResult:
    pkg = json.loads((ROOT / "package.json").read_text(encoding="utf-8"))
    declared = [cmd["command"] for cmd in pkg["contributes"]["commands"]]
    source = (ROOT / "src" / "extension.ts").read_text(encoding="utf-8")
    registered = re.findall(r'registerCommand\(\s*["\']([^"\']+)["\']', source)
    missing = sorted(set(declared) - set(registered))
    extra = sorted(set(registered) - set(declared))
    ok = not missing
    body = [
        f"declared_commands: {len(declared)}",
        f"registered_handlers: {len(registered)}",
        "missing_handlers: " + ("none" if not missing else ", ".join(missing)),
        "extra_internal_handlers: " + ("none" if not extra else ", ".join(extra)),
        "",
        "User functions covered:",
    ] + [f"- {item}" for item in declared]
    img = screenshot(
        "02_command_registration",
        "VS Code Command Registration",
        "package.json vs src/extension.ts",
        "PASS" if ok else "WARN",
        body[:20],
        "\n".join(body[20:]),
    )
    return CheckResult("Command registration", "static parity check", ok, body, "", img)


def parse_semgrep(label: str, command: list[str], image_name: str) -> CheckResult:
    result = run_cmd(label, command, timeout=240)
    ok = result.returncode == 0
    body: list[str] = [f"exit_code: {result.returncode}"]
    raw_tail = tail(result.stdout, result.stderr)
    try:
        data = json.loads(result.stdout)
        findings = data.get("results", [])
        ids = [item.get("check_id", "unknown") for item in findings]
        taint = [
            item
            for item in findings
            if item.get("extra", {}).get("dataflow_trace") or ".taint." in item.get("check_id", "")
        ]
        scanned = data.get("paths", {}).get("scanned", [])
        body.extend(
            [
                f"semgrep_version: {data.get('version', 'unknown')}",
                f"scanned_files: {len(scanned) or 'unknown'}",
                f"findings: {len(findings)}",
                f"taint_findings: {len(taint)}",
                "sample_rule_ids:",
            ]
        )
        body.extend([f"- {item}" for item in ids[:8]])
        raw_tail = "\n".join(
            [
                f"Command: {' '.join(command)}",
                f"First finding paths: {[item.get('path') for item in findings[:4]]}",
            ]
        )
    except Exception as exc:
        body.extend(["json_parse: failed", f"reason: {exc}"])
        ok = False
    img = screenshot(
        image_name,
        label,
        " ".join(command),
        "PASS" if ok else "WARN",
        body,
        raw_tail,
    )
    return CheckResult(label, " ".join(command), ok, body, raw_tail, img)


def build_policy_check() -> CheckResult:
    node = shutil.which("node") or "node"
    script = (
        "const path=require('path');"
        "const {loadPolicy}=require('./out/policy');"
        "const r=loadPolicy(path.join(process.cwd(),'test-samples'),process.cwd());"
        "console.log(JSON.stringify({ok:r.ok,presets:r.policy.presets,"
        "minSeverity:r.policy.severity.minSeverity,customRules:r.policy.rules.length,"
        "disabledRules:r.policy.disabledRules,excludes:r.policy.files.exclude},null,2));"
    )
    result = run_cmd("policy", [node, "-e", script], timeout=60)
    ok = result.returncode == 0
    body = [f"exit_code: {result.returncode}"]
    try:
        data = json.loads(result.stdout)
        ok = ok and bool(data.get("ok"))
        body.extend(
            [
                f"policy_ok: {data.get('ok')}",
                "active_presets: " + ", ".join(data.get("presets", [])),
                f"min_severity: {data.get('minSeverity')}",
                f"custom_rules_loaded: {data.get('customRules')}",
                f"disabled_rules: {len(data.get('disabledRules', []))}",
                f"exclude_globs: {len(data.get('excludes', []))}",
            ]
        )
    except Exception as exc:
        body.append(f"parse_failed: {exc}")
        ok = False
    img = screenshot(
        "03_policy_resolution",
        "Policy Resolution",
        "loadPolicy(test-samples)",
        "PASS" if ok else "WARN",
        body,
        result.stdout + result.stderr,
    )
    return CheckResult("Policy resolution", "loadPolicy(test-samples)", ok, body, tail(result.stdout, result.stderr), img)


def build_rules_index_check() -> CheckResult:
    node = shutil.which("node") or "node"
    script = (
        "const path=require('path');"
        "const {buildRulesIndex}=require('./out/rulesIndex');"
        "const i=buildRulesIndex(process.cwd(),path.join(process.cwd(),'test-samples'));"
        "console.log(JSON.stringify({files:i.files.map(f=>({path:f.path,source:f.source,enabled:f.enabled,ruleCount:f.ruleCount})),"
        "totalRules:i.rules.length,enabledRules:i.rules.filter(r=>r.enabled).length,"
        "taintRules:i.rules.filter(r=>r.mode==='taint').length},null,2));"
    )
    result = run_cmd("rules-index", [node, "-e", script], timeout=60)
    ok = result.returncode == 0
    body = [f"exit_code: {result.returncode}"]
    try:
        data = json.loads(result.stdout)
        body.extend(
            [
                f"policy_files_seen: {len(data.get('files', []))}",
                f"total_rules: {data.get('totalRules')}",
                f"enabled_rules: {data.get('enabledRules')}",
                f"taint_rules: {data.get('taintRules')}",
                "files:",
            ]
        )
        for item in data.get("files", []):
            body.append(f"- {item['path']} | {item['source']} | enabled={item['enabled']} | rules={item['ruleCount']}")
        ok = ok and data.get("totalRules", 0) > 0
    except Exception as exc:
        body.append(f"parse_failed: {exc}")
        ok = False
    img = screenshot(
        "06_rules_index",
        "Rules Index Verification",
        "buildRulesIndex(extensionRoot, test-samples)",
        "PASS" if ok else "WARN",
        body,
        result.stdout + result.stderr,
    )
    return CheckResult("Rules index", "buildRulesIndex", ok, body, tail(result.stdout, result.stderr), img)


def build_llm_prompt_check() -> CheckResult:
    node = shutil.which("node") or "node"
    script = (
        "const m=require('./out/llmModels');"
        "console.log(JSON.stringify({providers:m.ALL_PROVIDERS,defaults:m.PROVIDER_DEFAULT_MODEL,"
        "presetCounts:Object.fromEntries(Object.entries(m.PROVIDER_MODEL_PRESETS).map(([k,v])=>[k,v.length]))},null,2));"
    )
    result = run_cmd("llm-models", [node, "-e", script], timeout=60)
    ok = result.returncode == 0
    body = [f"exit_code: {result.returncode}", "prompt_commands: vuln, file, project, copy all"]
    try:
        data = json.loads(result.stdout)
        body.append("providers: " + ", ".join(data.get("providers", [])))
        for key, val in data.get("defaults", {}).items():
            body.append(f"- {key}: {val}")
        ok = ok and len(data.get("providers", [])) >= 4
    except Exception as exc:
        body.append(f"parse_failed: {exc}")
        ok = False
    img = screenshot(
        "07_llm_prompt_surface",
        "LLM Provider and Prompt Surface",
        "llmModels + prompt command surface",
        "PASS" if ok else "WARN",
        body,
        result.stdout + result.stderr,
    )
    return CheckResult("LLM prompt surface", "llmModels", ok, body, tail(result.stdout, result.stderr), img)


def build_webview_asset_check() -> CheckResult:
    files = [
        ROOT / "media" / "design" / "main.js",
        ROOT / "media" / "design" / "styles.css",
        ROOT / "media" / "design" / "controlCenter.js",
        ROOT / "media" / "design" / "controlCenter.css",
        ROOT / "src" / "panelMessages.ts",
        ROOT / "src" / "controlCenterMessages.ts",
    ]
    body = []
    ok = True
    for item in files:
        exists = item.exists()
        ok = ok and exists
        size = item.stat().st_size if exists else 0
        body.append(f"{item.relative_to(ROOT)} | exists={exists} | bytes={size}")
    img = screenshot(
        "08_webview_assets",
        "Webview Bundles and Message Contracts",
        "media/design + src/*Messages.ts",
        "PASS" if ok else "WARN",
        body,
        "",
    )
    return CheckResult("Webview assets", "asset existence check", ok, body, "", img)


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Segoe UI"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = NAVY
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.86), Inches(11.7), Inches(0.35))
        p2 = sub.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(11)
        p2.font.color.rgb = MUTED


def add_footer(slide, text: str = "VibeSec local functionality testing evidence") -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.12), Inches(12.2), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(8.5)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_text(slide, x, y, w, h, text, size=13, color=INK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def add_rect(slide, x, y, w, h, text, fill=WHITE, line=LINE, size=12, bold=False, color=INK):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_oval(slide, x, y, w, h, text, fill=WHITE, line=BLUE, size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Segoe UI"
    p.font.size = Pt(size)
    p.font.color.rgb = INK
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_line(slide, x1, y1, x2, y2, color=LINE, width=1.2):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def new_slide(prs: Presentation, title: str, subtitle: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    add_title(slide, title, subtitle)
    add_footer(slide)
    return slide


def add_image_slide(prs: Presentation, title: str, subtitle: str, image_path: str):
    slide = new_slide(prs, title, subtitle)
    slide.shapes.add_picture(image_path, Inches(0.75), Inches(1.35), width=Inches(11.85))
    return slide


def title_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, 0.85, 1.0, 11.7, 0.6, "VibeSec", size=36, color=WHITE, bold=True)
    add_text(
        slide,
        0.88,
        1.72,
        11.6,
        1.0,
        "Functionality Testing Evidence and UML Diagrams",
        size=30,
        color=RGBColor(224, 242, 254),
        bold=True,
    )
    add_text(
        slide,
        0.9,
        3.0,
        10.5,
        0.65,
        "Local VS Code extension verification deck covering the command surface, scan engine, policy/rules model, webviews, prompt surface, and architecture diagrams.",
        size=16,
        color=RGBColor(203, 213, 225),
    )
    add_text(slide, 0.9, 6.72, 7.5, 0.3, f"Generated from workspace: {ROOT}", size=8.5, color=RGBColor(148, 163, 184))


def summary_slide(prs: Presentation, checks: list[CheckResult]):
    slide = new_slide(prs, "Functional Testing Summary", "All evidence was generated from the local workspace.")
    passed = sum(1 for c in checks if c.ok)
    total = len(checks)
    add_rect(slide, 0.75, 1.35, 2.15, 1.0, f"{passed}/{total}\nchecks passed", fill=RGBColor(220, 252, 231), line=GREEN, size=18, bold=True, color=GREEN)
    add_rect(slide, 3.12, 1.35, 2.15, 1.0, "15\nVS Code commands", fill=RGBColor(219, 234, 254), line=BLUE, size=18, bold=True, color=BLUE)
    add_rect(slide, 5.49, 1.35, 2.15, 1.0, "154\nbundled rules", fill=RGBColor(224, 242, 254), line=CYAN, size=18, bold=True, color=CYAN)
    add_rect(slide, 7.86, 1.35, 2.15, 1.0, "13\ntaint rules", fill=RGBColor(237, 233, 254), line=PURPLE, size=18, bold=True, color=PURPLE)
    add_rect(slide, 10.23, 1.35, 2.15, 1.0, "React\nwebviews", fill=RGBColor(254, 243, 199), line=AMBER, size=18, bold=True, color=AMBER)

    rows = [["Area", "Evidence", "Status"]]
    for check in checks:
        rows.append([check.name, shorten("; ".join(check.summary[:3]), width=72, placeholder="..."), "PASS" if check.ok else "WARN"])
    table = slide.shapes.add_table(len(rows), 3, Inches(0.75), Inches(2.72), Inches(11.85), Inches(3.8)).table
    table.columns[0].width = Inches(2.45)
    table.columns[1].width = Inches(7.55)
    table.columns[2].width = Inches(1.35)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.size = Pt(9 if r else 10)
            p.font.bold = r == 0
            p.font.color.rgb = WHITE if r == 0 else INK
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = NAVY
            elif c == 2 and value == "PASS":
                cell.fill.fore_color.rgb = RGBColor(220, 252, 231)
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif c == 2:
                cell.fill.fore_color.rgb = RGBColor(254, 243, 199)
                p.font.color.rgb = AMBER
                p.font.bold = True
            else:
                cell.fill.fore_color.rgb = WHITE


def command_matrix_slide(prs: Presentation):
    pkg = json.loads((ROOT / "package.json").read_text(encoding="utf-8"))
    commands = [(cmd["command"], cmd.get("title", "")) for cmd in pkg["contributes"]["commands"]]
    groups = {
        "Scan": ["scanCurrentFile", "scanSelected", "scanWorkspace", "goToFinding"],
        "Policy": ["reloadPolicy", "openPolicyFile"],
        "Prompt": ["generatePrompts", "copyPromptForVuln", "copyPromptForFile", "copyPromptForAll"],
        "API": ["setApiKey", "clearApiKey", "testApiKey"],
        "UI": ["openControlCenter", "copyDescription"],
    }
    slide = new_slide(prs, "Function-by-Function Coverage", "Each contributed command is registered and covered by grouped evidence screenshots.")
    rows = [["Function", "User action", "Evidence"]]
    for command, title in commands:
        group = "General"
        for key, needles in groups.items():
            if any(needle in command for needle in needles):
                group = key
                break
        rows.append([command.replace("vibesec.", ""), title, f"{group} evidence + registration PASS"])

    table = slide.shapes.add_table(len(rows), 3, Inches(0.55), Inches(1.28), Inches(12.25), Inches(5.78)).table
    table.columns[0].width = Inches(3.05)
    table.columns[1].width = Inches(5.55)
    table.columns[2].width = Inches(3.65)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Segoe UI"
            p.font.size = Pt(8.3 if r else 9.5)
            p.font.bold = r == 0
            p.font.color.rgb = WHITE if r == 0 else INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE


def sequence_diagram_slide(prs: Presentation):
    slide = new_slide(prs, "Sequence Diagram", "Scan current file / selected files / workspace.")
    participants = [
        ("Developer", 0.55),
        ("VS Code\nCommand", 2.05),
        ("Extension\nOrchestrator", 3.75),
        ("Policy\nLoader", 5.55),
        ("Scanner", 7.2),
        ("Semgrep\nCLI", 8.72),
        ("Diagnostics\n+ Panel", 10.18),
        ("Prompt\nGenerator", 11.75),
    ]
    top, bottom = 1.35, 6.55
    for label, x in participants:
        add_rect(slide, x, top, 1.18, 0.55, label, fill=WHITE, line=BLUE, size=9.5, bold=True, color=NAVY)
        add_line(slide, x + 0.59, top + 0.62, x + 0.59, bottom, color=RGBColor(148, 163, 184), width=1)
    messages = [
        (0.95, 2.45, "click scan"),
        (2.45, 4.15, "execute command"),
        (4.15, 6.00, "load .vibesec.yaml"),
        (6.00, 4.15, "policy result"),
        (4.15, 7.65, "scan targets"),
        (7.65, 9.15, "semgrep scan --json"),
        (9.15, 7.65, "findings JSON"),
        (7.65, 10.75, "normalized findings"),
        (10.75, 4.15, "open finding / state"),
        (10.75, 12.25, "generate/copy fix prompt"),
    ]
    y = 2.2
    for x1, x2, label in messages:
        add_line(slide, x1, y, x2, y, color=BLUE if x2 > x1 else CYAN, width=1.4)
        add_text(slide, min(x1, x2) + 0.05, y - 0.22, abs(x2 - x1) - 0.05, 0.2, label, size=7.8, color=INK, align=PP_ALIGN.CENTER)
        y += 0.42


def use_case_diagram_slide(prs: Presentation):
    slide = new_slide(prs, "Use Case Diagram", "Developer-local security workflow.")
    add_rect(slide, 0.55, 3.05, 1.45, 0.65, "Developer", fill=RGBColor(219, 234, 254), line=BLUE, size=13, bold=True, color=BLUE)
    add_rect(slide, 10.95, 1.15, 1.75, 0.55, "Semgrep CLI", fill=RGBColor(240, 253, 244), line=GREEN, size=11, bold=True, color=GREEN)
    add_rect(slide, 10.95, 5.7, 1.75, 0.55, "LLM Provider", fill=RGBColor(237, 233, 254), line=PURPLE, size=11, bold=True, color=PURPLE)

    cases = [
        ("Scan current file", 2.55, 1.25),
        ("Scan selected files/folder", 4.95, 1.25),
        ("Scan workspace", 7.55, 1.25),
        ("View findings", 2.55, 2.45),
        ("Inspect taint flow", 4.95, 2.45),
        ("Go to code location", 7.55, 2.45),
        ("Manage policies", 2.55, 3.65),
        ("Enable/disable rules", 4.95, 3.65),
        ("Import/create rule file", 7.55, 3.65),
        ("Generate fix prompt", 2.55, 4.85),
        ("Manage API keys", 4.95, 4.85),
        ("View logs/history", 7.55, 4.85),
    ]
    for label, x, y in cases:
        add_oval(slide, x, y, 1.85, 0.62, label, fill=WHITE, line=BLUE, size=8.8)
        add_line(slide, 2.0, 3.38, x + 0.02, y + 0.31, color=RGBColor(148, 163, 184), width=0.8)
    for x, y in [(2.55, 1.25), (4.95, 1.25), (7.55, 1.25)]:
        add_line(slide, x + 1.85, y + 0.31, 10.95, 1.42, color=GREEN, width=1)
    for x, y in [(2.55, 4.85), (4.95, 4.85)]:
        add_line(slide, x + 1.85, y + 0.31, 10.95, 5.97, color=PURPLE, width=1)


def class_diagram_slide(prs: Presentation):
    slide = new_slide(prs, "Class Diagram", "Main TypeScript modules and responsibilities.")
    boxes = [
        ("extension.ts\n+ activate()\n+ runScanOnTargets()\n+ command handlers", 0.55, 1.25, BLUE),
        ("FindingsProvider\n+ PanelState\n+ prompt cache\n+ tree data", 3.05, 1.25, CYAN),
        ("PanelController\n+ webview HTML\n+ message router\n+ file tree", 5.55, 1.25, CYAN),
        ("ControlCenterController\n+ settings\n+ policies/rules\n+ logs/history", 8.05, 1.25, PURPLE),
        ("policy.ts\n+ loadPolicy()\n+ activePolicyFiles\n+ validation/merge", 0.55, 3.25, GREEN),
        ("scanner.ts\n+ runSemgrep()\n+ parseFinding()\n+ taint trace", 3.05, 3.25, GREEN),
        ("rulesIndex.ts\n+ bundled/custom files\n+ disabled rules\n+ rule metadata", 5.55, 3.25, GREEN),
        ("promptGenerator.ts\n+ vuln/file/project prompts\n+ source context", 8.05, 3.25, AMBER),
        ("llmClient.ts\n+ provider clients\n+ timeout/errors\n+ custom endpoint", 10.55, 3.25, AMBER),
        ("LogBus / LogStore\n+ ring buffer\n+ JSONL persistence\n+ output channel", 1.8, 5.25, RED),
        ("ScanHistoryStore\n+ workspaceState\n+ recent scans", 4.3, 5.25, RED),
        ("types.ts\n+ Finding\n+ PolicyConfig\n+ CustomRule", 6.8, 5.25, NAVY),
        ("secrets.ts\n+ SecretStorage\n+ provider key flow", 9.3, 5.25, PURPLE),
    ]
    for label, x, y, color in boxes:
        add_rect(slide, x, y, 2.0, 1.0, label, fill=WHITE, line=color, size=8.5, bold=False, color=INK)
    lines = [
        (1.55, 2.25, 1.55, 3.25),
        (1.55, 2.25, 4.05, 3.25),
        (1.55, 2.25, 6.55, 3.25),
        (1.55, 2.25, 9.05, 3.25),
        (1.55, 2.25, 4.05, 1.25),
        (4.05, 2.25, 6.55, 1.25),
        (9.05, 2.25, 11.55, 3.25),
        (9.05, 4.25, 11.55, 4.25),
        (4.05, 4.25, 2.8, 5.25),
        (1.55, 2.25, 5.3, 5.25),
        (4.05, 4.25, 7.8, 5.25),
    ]
    for line in lines:
        add_line(slide, *line, color=RGBColor(100, 116, 139), width=1)


def next_steps_slide(prs: Presentation):
    slide = new_slide(prs, "Testing Notes", "What this deck proves and what still needs interactive VS Code validation.")
    add_rect(slide, 0.8, 1.35, 5.55, 4.7, "", fill=WHITE, line=LINE)
    add_text(slide, 1.05, 1.62, 5.05, 0.3, "Verified locally", size=16, color=GREEN, bold=True)
    verified = [
        "Build and bundled webview assets compile.",
        "All contributed VS Code commands have registered handlers.",
        "Policy loader resolves the sample multi-policy config.",
        "Semgrep normal scan returns findings from sample JS.",
        "Semgrep taint scan returns normal and taint findings.",
        "Rules index lists bundled/default/taint/custom policy files.",
        "LLM provider/model surface is wired for prompt commands.",
    ]
    y = 2.08
    for item in verified:
        add_text(slide, 1.05, y, 4.95, 0.28, "- " + item, size=11, color=INK)
        y += 0.43

    add_rect(slide, 6.95, 1.35, 5.55, 4.7, "", fill=WHITE, line=LINE)
    add_text(slide, 7.2, 1.62, 5.05, 0.3, "Recommended manual pass", size=16, color=AMBER, bold=True)
    manual = [
        "Launch Extension Development Host.",
        "Click each command from Command Palette.",
        "Capture real VS Code UI screenshots if needed for submission.",
        "Test API calls only with valid provider keys.",
        "Try policy create/import/delete from Control Center.",
    ]
    y = 2.08
    for item in manual:
        add_text(slide, 7.2, y, 4.95, 0.28, "- " + item, size=11, color=INK)
        y += 0.43


def create_deck(checks: list[CheckResult]) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    title_slide(prs)
    summary_slide(prs, checks)
    command_matrix_slide(prs)
    for check in checks:
        if check.image:
            add_image_slide(prs, "Evidence Screenshot: " + check.name, check.command, check.image)
    sequence_diagram_slide(prs)
    use_case_diagram_slide(prs)
    class_diagram_slide(prs)
    next_steps_slide(prs)
    PPTX_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_OUT)


def main() -> None:
    ASSETS.mkdir(parents=True, exist_ok=True)
    checks: list[CheckResult] = []
    checks.append(build_compile_check())
    checks.append(build_command_registration_check())
    checks.append(build_policy_check())

    semgrep = shutil.which("semgrep") or "semgrep"
    checks.append(
        parse_semgrep(
            "Normal Semgrep Scan",
            [semgrep, "scan", "--json", "--metrics=off", "--config", "rules/default.yaml", "test-samples/simple_web_policy_test.js"],
            "04_normal_scan",
        )
    )
    checks.append(
        parse_semgrep(
            "Normal + Taint Semgrep Scan",
            [
                semgrep,
                "scan",
                "--json",
                "--metrics=off",
                "--config",
                "rules/default.yaml",
                "--config",
                "rules/taint.yaml",
                "test-samples/vibesec_11_normal_3_taint.py",
            ],
            "05_taint_scan",
        )
    )
    checks.append(build_rules_index_check())
    checks.append(build_llm_prompt_check())
    checks.append(build_webview_asset_check())

    report = {
        "workspace": str(ROOT),
        "pptx": str(PPTX_OUT),
        "checks": [
            {
                "name": check.name,
                "command": check.command,
                "ok": check.ok,
                "summary": check.summary,
                "image": check.image,
            }
            for check in checks
        ],
    }
    REPORT_OUT.write_text(json.dumps(report, indent=2), encoding="utf-8")
    create_deck(checks)
    print(PPTX_OUT)


if __name__ == "__main__":
    os.chdir(ROOT)
    main()
